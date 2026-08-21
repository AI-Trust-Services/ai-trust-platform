"""Document parsing for AI-assisted registration.

Dispatches on file extension: PDF (pypdf), DOCX (python-docx), PPTX (python-pptx),
TXT/MD (decode), images (Pillow → resize longest side ≤1568px → base64 for vision).
Text is truncated to ~15 000 chars. Returns a ParsedDoc carrying either extracted
text or a base64 image payload, plus an ``is_image`` flag so the router picks
LLM_MODEL vs LLM_VISION_MODEL.
"""
from __future__ import annotations

import base64
import io
import os
from dataclasses import dataclass

from PIL import UnidentifiedImageError as _PILUnidentifiedImageError

from ai_trust_logging import get_logger

logger = get_logger(__name__)

TEXT_EXTENSIONS = {".txt", ".md", ".markdown"}
PDF_EXTENSIONS = {".pdf"}
WORD_EXTENSIONS = {".docx"}  # python-docx only supports OOXML (.docx); legacy .doc is not supported
PPTX_EXTENSIONS = {".ppt", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
ALL_SUPPORTED = TEXT_EXTENSIONS | PDF_EXTENSIONS | WORD_EXTENSIONS | PPTX_EXTENSIONS | IMAGE_EXTENSIONS

MAX_TEXT_LENGTH = int(os.environ.get("ASSIST_MAX_TEXT_LENGTH", "15000"))
MAX_IMAGE_SIDE = 1568  # vision-model recommended long side
MAX_FILE_BYTES = 20 * 1024 * 1024  # 20 MB application-level cap (nginx allows 50 MB)

_MEDIA_TYPES = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
}


class DocumentParseError(Exception):
    """Raised when a document cannot be parsed."""


@dataclass
class ParsedDoc:
    filename: str
    text: str | None = None
    image_b64: str | None = None
    media_type: str | None = None

    @property
    def is_image(self) -> bool:
        return self.image_b64 is not None


def _ext(filename: str) -> str:
    return os.path.splitext(filename)[1].lower()


def is_supported(filename: str) -> bool:
    return _ext(filename) in ALL_SUPPORTED


def _parse_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _parse_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def _parse_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if len(slide_texts) > 1:
            parts.append("\n".join(slide_texts))
    return "\n\n".join(parts)


def _parse_image(content: bytes, ext: str) -> tuple[str, str]:
    media_type = _MEDIA_TYPES.get(ext, "image/png")
    try:
        from PIL import Image

        img = Image.open(io.BytesIO(content))
        if max(img.size) > MAX_IMAGE_SIDE:
            ratio = MAX_IMAGE_SIDE / max(img.size)
            img = img.resize((int(img.size[0] * ratio), int(img.size[1] * ratio)), Image.Resampling.LANCZOS)
            buffer = io.BytesIO()
            img.save(buffer, format="PNG" if ext == ".png" else "JPEG")
            content = buffer.getvalue()
            logger.info("document.image_resized", extra={"new_size": img.size})
    except (OSError, _PILUnidentifiedImageError) as exc:  # resize is best-effort; fall back to original bytes
        logger.warning("document.image_resize_failed", extra={"error": str(exc)})
    return base64.standard_b64encode(content).decode("utf-8"), media_type


def parse_document(filename: str, content: bytes) -> ParsedDoc:
    """Parse ``content`` by extension. Raises DocumentParseError on unsupported/failed parse."""
    if len(content) > MAX_FILE_BYTES:
        raise DocumentParseError(
            f"File too large: {len(content) / (1024 * 1024):.1f} MB (max {MAX_FILE_BYTES // (1024 * 1024)} MB)"
        )

    ext = _ext(filename)
    if ext not in ALL_SUPPORTED:
        raise DocumentParseError(
            f"Unsupported file type: {ext}. Supported: {', '.join(sorted(ALL_SUPPORTED))}"
        )

    try:
        if ext in IMAGE_EXTENSIONS:
            b64, media_type = _parse_image(content, ext)
            return ParsedDoc(filename=filename, image_b64=b64, media_type=media_type)

        if ext in TEXT_EXTENSIONS:
            text = content.decode("utf-8", errors="replace")
        elif ext in PDF_EXTENSIONS:
            text = _parse_pdf(content)
        elif ext in WORD_EXTENSIONS:
            text = _parse_docx(content)
        else:  # PPTX_EXTENSIONS
            text = _parse_pptx(content)
    except DocumentParseError:
        raise
    except Exception as exc:
        logger.warning("document.parse_failed", extra={"file_name": filename, "error": str(exc)})
        raise DocumentParseError(f"Failed to parse {ext} document: {exc}") from exc

    if len(text) > MAX_TEXT_LENGTH:
        text = text[:MAX_TEXT_LENGTH] + "\n\n[Document truncated due to length…]"
        logger.info("document.text_truncated", extra={"file_name": filename})
    return ParsedDoc(filename=filename, text=text)
